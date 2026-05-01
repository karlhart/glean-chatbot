"""Generate a 4-slide PowerPoint presentation for the Lumina Glean Chatbot project."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from pptx.enum.dml import MSO_THEME_COLOR
import copy

# ── Brand colours ────────────────────────────────────────────────────────────
GLEAN_BLUE    = RGBColor(0x17, 0x3B, 0xD6)   # Glean brand blue
GLEAN_DARK    = RGBColor(0x0D, 0x1B, 0x3E)   # dark navy
ACCENT_GOLD   = RGBColor(0xF5, 0xA6, 0x23)   # warm gold
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY    = RGBColor(0xF2, 0xF4, 0xF8)
MID_GREY      = RGBColor(0x8A, 0x94, 0xA6)
BOX_GREEN     = RGBColor(0x1A, 0x8C, 0x5A)
BOX_BLUE      = RGBColor(0x17, 0x3B, 0xD6)
BOX_ORANGE    = RGBColor(0xE8, 0x6B, 0x1F)
BOX_PURPLE    = RGBColor(0x6C, 0x3F, 0xAF)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


# ── Helpers ───────────────────────────────────────────────────────────────────

def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]          # truly blank layout
    return prs.slides.add_slide(blank_layout)


def fill_bg(slide, color: RGBColor):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()   # no border
    return shape


def add_text_box(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=WHITE,
                 align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox


def add_label_box(slide, label, left, top, width, height,
                  bg_color, text_color=WHITE, font_size=14, bold=True):
    """Rounded-look box: filled rect + centred text."""
    add_rect(slide, left, top, width, height, bg_color)
    add_text_box(slide, label, left, top, width, height,
                 font_size=font_size, bold=bold, color=text_color,
                 align=PP_ALIGN.CENTER)


def add_arrow(slide, x1, y1, x2, y2, color=MID_GREY, width_pt=2):
    """Draw a connector line with an arrowhead."""
    connector = slide.shapes.add_connector(
        1,   # STRAIGHT connector
        x1, y1, x2, y2
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(width_pt)
    # add arrowhead at end
    cNvSpPr = connector._element.find('.//' + qn('p:nvCxnSpPr'))
    ln = connector.line._ln
    tailEnd = ln.find(qn('a:tailEnd'))
    if tailEnd is None:
        from lxml import etree
        ln.append(etree.fromstring(
            '<a:tailEnd xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
            'type="none"/>'
        ))
    from lxml import etree
    # set head end arrow
    headEnd = ln.find(qn('a:headEnd'))
    if headEnd is not None:
        ln.remove(headEnd)
    ln.append(etree.fromstring(
        '<a:headEnd xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
        'type="arrow" w="med" len="med"/>'
    ))
    return connector


# ── Slide builders ────────────────────────────────────────────────────────────

def slide_1_company(prs):
    """Background on Lumina Stream Studios."""
    slide = blank_slide(prs)
    fill_bg(slide, GLEAN_DARK)

    # left navy panel
    add_rect(slide, 0, 0, Inches(5), SLIDE_H, RGBColor(0x0A, 0x12, 0x2E))

    # slide number pill
    add_rect(slide, Inches(0.3), Inches(0.3), Inches(0.45), Inches(0.28), GLEAN_BLUE)
    add_text_box(slide, "01", Inches(0.3), Inches(0.28), Inches(0.45), Inches(0.32),
                 font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # company name — left panel
    add_text_box(slide, "LUMINA\nSTREAM\nSTUDIOS",
                 Inches(0.35), Inches(0.75), Inches(4.3), Inches(2.5),
                 font_size=36, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # gold underline bar
    add_rect(slide, Inches(0.35), Inches(3.35), Inches(1.2), Inches(0.06), ACCENT_GOLD)

    # tagline
    add_text_box(slide, "Global Indie-Major Production House\n& Streaming Platform",
                 Inches(0.35), Inches(3.5), Inches(4.3), Inches(0.8),
                 font_size=13, bold=False, color=MID_GREY, align=PP_ALIGN.LEFT)

    # key stats row
    stats = [("14", "Countries"), ("7", "Indexed\nDocuments"), ("3", "Glean APIs\nUsed")]
    for i, (num, lbl) in enumerate(stats):
        x = Inches(0.35 + i * 1.55)
        add_text_box(slide, num, x, Inches(4.5), Inches(1.4), Inches(0.6),
                     font_size=32, bold=True, color=ACCENT_GOLD, align=PP_ALIGN.LEFT)
        add_text_box(slide, lbl, x, Inches(5.1), Inches(1.4), Inches(0.6),
                     font_size=10, color=MID_GREY, align=PP_ALIGN.LEFT)

    # ── right panel ──────────────────────────────────────────────────────────
    # section header
    add_text_box(slide, "THE CHALLENGE",
                 Inches(5.4), Inches(0.55), Inches(7.5), Inches(0.4),
                 font_size=11, bold=True, color=ACCENT_GOLD, align=PP_ALIGN.LEFT)
    add_text_box(slide, "Knowledge scattered across isolated data islands",
                 Inches(5.4), Inches(0.9), Inches(7.5), Inches(0.5),
                 font_size=20, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # data island cards
    islands = [
        ("📁  Google Drive",    "Scripts & production\nschedules",          BOX_GREEN),
        ("🗂  Box",             "Contracts & talent\nbuyout agreements",     BOX_BLUE),
        ("💬  Slack + Outlook", "1,000+ channels &\nemail threads",         BOX_ORANGE),
        ("🔧  Jira + Confluence","Post-production workflows\n& VFX tracking", BOX_PURPLE),
    ]
    cols = 2
    card_w, card_h = Inches(3.4), Inches(1.35)
    x_starts = [Inches(5.4), Inches(9.1)]
    y_starts = [Inches(1.65), Inches(3.2)]
    for idx, (title, desc, color) in enumerate(islands):
        col = idx % cols
        row = idx // cols
        x = x_starts[col]
        y = y_starts[row]
        add_rect(slide, x, y, card_w, card_h, color)
        add_text_box(slide, title, x + Inches(0.15), y + Inches(0.12),
                     card_w - Inches(0.3), Inches(0.38),
                     font_size=13, bold=True, color=WHITE)
        add_text_box(slide, desc, x + Inches(0.15), y + Inches(0.5),
                     card_w - Inches(0.3), Inches(0.72),
                     font_size=11, color=WHITE)

    # solution statement
    add_rect(slide, Inches(5.4), Inches(4.75), Inches(7.5), Inches(0.06), ACCENT_GOLD)
    add_text_box(slide,
                 "Solution: Index all internal documents into Glean → "
                 "single AI-powered knowledge base for every Lumina employee.",
                 Inches(5.4), Inches(4.9), Inches(7.5), Inches(0.8),
                 font_size=13, italic=True, color=WHITE)

    # source docs list
    add_text_box(slide, "Documents Indexed:",
                 Inches(5.4), Inches(5.85), Inches(3.5), Inches(0.3),
                 font_size=10, bold=True, color=ACCENT_GOLD)
    docs = ("Employee Onboarding  •  IT Security  •  Production Workflow  •  "
            "Legal & Contracts  •  Post-Production / VFX  •  "
            "International Co-Production  •  Content Delivery Standards")
    add_text_box(slide, docs,
                 Inches(5.4), Inches(6.15), Inches(7.5), Inches(1.0),
                 font_size=10, color=MID_GREY)


def slide_2_requirements(prs):
    """Project requirements."""
    slide = blank_slide(prs)
    fill_bg(slide, RGBColor(0xF8, 0xF9, 0xFB))

    # top bar
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.1), GLEAN_DARK)
    add_text_box(slide, "02", Inches(0.35), Inches(0.3), Inches(0.4), Inches(0.5),
                 font_size=10, bold=True, color=ACCENT_GOLD, align=PP_ALIGN.CENTER)
    add_text_box(slide, "PROJECT REQUIREMENTS",
                 Inches(0.85), Inches(0.32), Inches(6), Inches(0.5),
                 font_size=22, bold=True, color=WHITE)
    add_text_box(slide, "Glean SA Technical Exercise",
                 Inches(0.85), Inches(0.72), Inches(6), Inches(0.35),
                 font_size=12, color=MID_GREY)

    # ── Three API requirement boxes ──────────────────────────────────────────
    api_cards = [
        ("INDEXING API",
         "Push internal documents\ninto a custom datasource",
         "• Stable document IDs\n• Custom objectType\n• Permission ACLs\n• viewURL matching urlRegex",
         GLEAN_BLUE),
        ("SEARCH API",
         "Retrieve relevant documents\nfor a natural-language query",
         "• Keyword extraction\n• datasourcesFilter scoping\n• Top-K result retrieval\n• Snippet enrichment",
         BOX_PURPLE),
        ("CHAT API",
         "Generate a grounded answer\nwith source citations",
         "• Context injection (RAG)\n• Anti-hallucination prompt\n• [1] [2] citation format\n• 25s timeout + fallback",
         BOX_GREEN),
    ]
    card_w = Inches(3.9)
    for i, (title, sub, bullets, color) in enumerate(api_cards):
        x = Inches(0.3 + i * 4.35)
        y = Inches(1.3)
        add_rect(slide, x, y, card_w, Inches(0.55), color)
        add_text_box(slide, title, x + Inches(0.15), y + Inches(0.08),
                     card_w, Inches(0.4), font_size=13, bold=True, color=WHITE)
        add_rect(slide, x, y + Inches(0.55), card_w, Inches(2.6),
                 RGBColor(0xFF, 0xFF, 0xFF), color)
        add_text_box(slide, sub, x + Inches(0.15), y + Inches(0.65),
                     card_w - Inches(0.3), Inches(0.7),
                     font_size=12, bold=True, color=GLEAN_DARK)
        add_text_box(slide, bullets, x + Inches(0.15), y + Inches(1.35),
                     card_w - Inches(0.3), Inches(1.7),
                     font_size=11, color=RGBColor(0x44, 0x44, 0x55))

    # ── MCP requirement ──────────────────────────────────────────────────────
    add_rect(slide, Inches(0.3), Inches(4.15), Inches(12.7), Inches(0.5), GLEAN_DARK)
    add_text_box(slide, "MCP TOOL REQUIREMENT",
                 Inches(0.5), Inches(4.22), Inches(4), Inches(0.36),
                 font_size=13, bold=True, color=ACCENT_GOLD)
    add_text_box(slide,
                 "Expose the full pipeline as a single MCP tool callable from Claude Desktop or Cursor",
                 Inches(4.5), Inches(4.22), Inches(8.2), Inches(0.36),
                 font_size=13, color=WHITE)

    mcp_items = [
        ("ask_lumina( question )",         "Required: natural-language question",                       GLEAN_BLUE),
        ("datasource, top_k",              "Optional: scope and result count",                           BOX_PURPLE),
        ("after_date, before_date",        "Optional: date range filters",                               BOX_GREEN),
        ("fast_mode",                      "Optional: skip Chat for ~800ms response (vs 10–20s)",       BOX_ORANGE),
    ]
    for i, (param, desc, color) in enumerate(mcp_items):
        x = Inches(0.3)
        y = Inches(4.85 + i * 0.52)
        add_rect(slide, x, y, Inches(3.2), Inches(0.38), color)
        add_text_box(slide, param, x + Inches(0.12), y + Inches(0.05),
                     Inches(3.0), Inches(0.3), font_size=11, bold=True, color=WHITE)
        add_text_box(slide, desc, Inches(3.65), y + Inches(0.05),
                     Inches(9.2), Inches(0.3), font_size=11,
                     color=RGBColor(0x33, 0x33, 0x44))


def slide_3_architecture(prs):
    """Architecture overview with flow diagram."""
    slide = blank_slide(prs)
    fill_bg(slide, GLEAN_DARK)

    # header
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.95), RGBColor(0x0A, 0x12, 0x2E))
    add_text_box(slide, "03", Inches(0.35), Inches(0.22), Inches(0.4), Inches(0.5),
                 font_size=10, bold=True, color=ACCENT_GOLD, align=PP_ALIGN.CENTER)
    add_text_box(slide, "ARCHITECTURE OVERVIEW",
                 Inches(0.85), Inches(0.18), Inches(8), Inches(0.45),
                 font_size=22, bold=True, color=WHITE)
    add_text_box(slide, "RAG pipeline: Index → Search → Enrich → Chat → Answer + Sources",
                 Inches(0.85), Inches(0.6), Inches(11), Inches(0.32),
                 font_size=12, color=MID_GREY)

    # ── INGEST FLOW (top row) ────────────────────────────────────────────────
    add_text_box(slide, "INGEST  (one-time / on update)",
                 Inches(0.35), Inches(1.1), Inches(8), Inches(0.3),
                 font_size=10, bold=True, color=ACCENT_GOLD)

    ingest_steps = [
        ("docs/*.md\n7 Lumina\nDocuments",  RGBColor(0x2A, 0x3D, 0x6A)),
        ("indexer.py\nGlean Indexing API\n/indexdocuments", GLEAN_BLUE),
        ("interviewds\nDatasource\n(~15–20 min lag)", RGBColor(0x2A, 0x3D, 0x6A)),
    ]
    box_w, box_h = Inches(2.8), Inches(1.1)
    x_positions = [Inches(0.35), Inches(3.8), Inches(7.25)]
    for i, (label, color) in enumerate(ingest_steps):
        x = x_positions[i]
        y = Inches(1.45)
        add_rect(slide, x, y, box_w, box_h, color)
        add_text_box(slide, label, x + Inches(0.1), y + Inches(0.08),
                     box_w - Inches(0.2), box_h - Inches(0.15),
                     font_size=11, bold=False, color=WHITE, align=PP_ALIGN.CENTER)

    # arrows between ingest boxes
    for i in range(len(x_positions) - 1):
        x1 = x_positions[i] + box_w
        x2 = x_positions[i + 1]
        mid_y = Inches(1.45) + box_h / 2
        add_arrow(slide, x1, mid_y, x2, mid_y, ACCENT_GOLD, 2)

    # ── QUERY FLOW (main diagram) ─────────────────────────────────────────────
    add_text_box(slide, "QUERY  (per ask_lumina invocation)",
                 Inches(0.35), Inches(2.75), Inches(8), Inches(0.3),
                 font_size=10, bold=True, color=ACCENT_GOLD)

    query_steps = [
        ("User\nQuestion",       RGBColor(0x33, 0x44, 0x66),  Inches(0.35)),
        ("Keyword\nExtraction",  RGBColor(0x1E, 0x3A, 0x8A),  Inches(2.2)),
        ("Glean\nSearch API",    GLEAN_BLUE,                   Inches(4.05)),
        ("Content\nEnrichment",  RGBColor(0x1E, 0x3A, 0x8A),  Inches(5.9)),
        ("Glean\nChat API",      BOX_PURPLE,                   Inches(7.75)),
        ("Answer +\nSources",    BOX_GREEN,                    Inches(9.6)),
    ]
    qbox_w, qbox_h = Inches(1.65), Inches(1.25)
    qbox_y = Inches(3.1)
    for label, color, x in query_steps:
        add_rect(slide, x, qbox_y, qbox_w, qbox_h, color)
        add_text_box(slide, label, x + Inches(0.08), qbox_y + Inches(0.15),
                     qbox_w - Inches(0.15), qbox_h - Inches(0.25),
                     font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # arrows between query boxes
    for i in range(len(query_steps) - 1):
        x1 = query_steps[i][2] + qbox_w
        x2 = query_steps[i + 1][2]
        mid_y = qbox_y + qbox_h / 2
        add_arrow(slide, x1, mid_y, x2, mid_y, ACCENT_GOLD, 2)

    # ── Annotation callouts ──────────────────────────────────────────────────
    # MCP wrapper bracket label
    add_rect(slide, Inches(0.35), Inches(4.5), Inches(10.9), Inches(0.04),
             RGBColor(0x55, 0x66, 0x88))
    add_text_box(slide, "⟵  Wrapped as ask_lumina MCP Tool  ⟶",
                 Inches(3.0), Inches(4.55), Inches(5.5), Inches(0.3),
                 font_size=10, italic=True, color=MID_GREY, align=PP_ALIGN.CENTER)

    # fast_mode callout
    add_rect(slide, Inches(7.75), Inches(4.9), Inches(3.5), Inches(0.6),
             RGBColor(0xE8, 0x6B, 0x1F))
    add_text_box(slide,
                 "fast_mode=True skips Chat\n→ ~800ms instead of 10–20s",
                 Inches(7.85), Inches(4.92), Inches(3.3), Inches(0.55),
                 font_size=10, color=WHITE)

    # ── Key design decisions ─────────────────────────────────────────────────
    add_text_box(slide, "KEY DESIGN DECISIONS",
                 Inches(0.35), Inches(5.65), Inches(6), Inches(0.28),
                 font_size=10, bold=True, color=ACCENT_GOLD)

    decisions = [
        "Official glean-api-client — correct datasourcesFilter + managed auth",
        "Fetch 20 results → URL allowlist (shared sandbox has other candidates' docs)",
        "Snippet-first truncation — Glean's top excerpt always included, never truncated away",
        "Anti-hallucination prompt — Chat instructed to cite sources or say 'I don't know'",
        "X-Glean-ActAs header — enforces per-user permissions with Global token",
    ]
    for i, d in enumerate(decisions):
        add_text_box(slide, f"▸  {d}",
                     Inches(0.35), Inches(5.95 + i * 0.28),
                     Inches(12.6), Inches(0.28),
                     font_size=10.5, color=WHITE)


def fixes_header(slide, slide_num, subtitle):
    """Shared header for the two fix slides."""
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.0), GLEAN_DARK)
    add_rect(slide, 0, 0, Inches(0.18), Inches(1.0), ACCENT_GOLD)
    add_text_box(slide, f"0{slide_num}",
                 Inches(0.3), Inches(0.25), Inches(0.45), Inches(0.45),
                 font_size=11, bold=True, color=ACCENT_GOLD, align=PP_ALIGN.CENTER)
    add_text_box(slide, "TROUBLESHOOTING & FIXES",
                 Inches(0.85), Inches(0.14), Inches(9), Inches(0.42),
                 font_size=22, bold=True, color=WHITE)
    add_text_box(slide, subtitle,
                 Inches(0.85), Inches(0.58), Inches(11), Inches(0.32),
                 font_size=12, color=MID_GREY)


def fix_row(slide, y, num, category, issue, fix, color):
    """Render a single fix row."""
    row_h = Inches(0.72)
    # number pill
    add_rect(slide, Inches(0.25), y + Inches(0.12), Inches(0.38), Inches(0.38), color)
    add_text_box(slide, str(num),
                 Inches(0.25), y + Inches(0.12), Inches(0.38), Inches(0.38),
                 font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # category badge
    add_rect(slide, Inches(0.75), y + Inches(0.15), Inches(1.85), Inches(0.26), color)
    add_text_box(slide, category,
                 Inches(0.75), y + Inches(0.15), Inches(1.85), Inches(0.26),
                 font_size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # issue
    add_text_box(slide, issue,
                 Inches(2.75), y + Inches(0.06), Inches(4.8), Inches(0.32),
                 font_size=11, bold=True, color=GLEAN_DARK)
    # fix
    add_text_box(slide, fix,
                 Inches(2.75), y + Inches(0.38), Inches(4.8), Inches(0.28),
                 font_size=10, color=RGBColor(0x44, 0x55, 0x66))
    # divider
    add_rect(slide, Inches(0.25), y + row_h - Inches(0.03),
             Inches(12.8), Inches(0.02), LIGHT_GREY)


def slide_4_fixes_a(prs):
    """Fixes 1–7."""
    slide = blank_slide(prs)
    fill_bg(slide, RGBColor(0xF8, 0xF9, 0xFB))
    fixes_header(slide, 4, "Fixes 1–7  ·  Root causes & resolutions from development")

    fixes = [
        (1, "INDEXING",     "HTTP 400: viewURL didn't match datasource urlRegex",
                            "Aligned INTRANET_BASE to the sandbox-configured regex pattern",
                            GLEAN_BLUE),
        (2, "RESILIENCE",   "No retry logic, logging, or env-var validation",
                            "Added exponential backoff on 429, structured logging with requestId, fail-fast startup check",
                            BOX_GREEN),
        (3, "BEST PRACTICE","Gaps vs Glean MCP server guidelines",
                            "Added keyword extraction, read_document enrichment pattern, anti-hallucination prompt, date filters",
                            BOX_PURPLE),
        (4, "PERFORMANCE",  "Chat API timeouts — 30KB prompts from full document content",
                            "Capped at 1500 chars/doc; snippet-first truncation ensures key passage is never cut",
                            BOX_ORANGE),
        (5, "ENRICHMENT",   "Enriched 0/5 — startswith() failed when Glean wraps viewURLs",
                            "Switched to substring match (INTRANET_BASE in url) to handle any redirect prefix",
                            GLEAN_BLUE),
        (6, "MCP TIMEOUT",  "Claude Desktop MCP call timing out — Chat taking >45s",
                            "Reduced timeout to 25s; added snippet fallback so tool always returns grounded content",
                            BOX_ORANGE),
        (7, "FILTERING",    "datasourcesFilter silently ignored by the REST API",
                            "URL post-filter on returned results — REST API field ignored at both top-level and requestOptions",
                            BOX_GREEN),
    ]

    for i, (num, cat, issue, fix, color) in enumerate(fixes):
        y = Inches(1.1 + i * 0.76)
        fix_row(slide, y, num, cat, issue, fix, color)

    # column headers
    for label, x in [("  #  CAT", Inches(0.25)), ("ISSUE", Inches(2.75)), ("FIX", Inches(7.7))]:
        add_text_box(slide, label, x, Inches(1.05), Inches(4.5), Inches(0.22),
                     font_size=8, bold=True, color=MID_GREY)


def slide_5_fixes_b(prs):
    """Fixes 8–13 + summary row."""
    slide = blank_slide(prs)
    fill_bg(slide, RGBColor(0xF8, 0xF9, 0xFB))
    fixes_header(slide, 5, "Fixes 8–13  ·  Root causes & resolutions from development")

    fixes = [
        (8,  "API CLIENT",  "Raw requests fragile — datasourcesFilter wrong, 5s default timeout",
                            "Switched to official glean-api-client; timeout_ms=25000 on client constructor",
                            BOX_PURPLE),
        (9,  "SANDBOX",     "Other candidates' documents returned (shared interviewds datasource)",
                            "Replaced prefix filter with exact URL allowlist of our 7 known document URLs",
                            GLEAN_BLUE),
        (10, "CITATIONS",   "Sources dropped — Claude Desktop paraphrases tool output",
                            "Added MCP server instructions field telling Claude to present response verbatim",
                            BOX_GREEN),
        (11, "KEYWORDS",    "Commas and punctuation leaking into search queries",
                            "Regex strip of all punctuation before tokenising keyword extractor",
                            BOX_ORANGE),
        (12, "SEARCH",      "Generic queries (e.g. 'summarize checklist') returned 0 Lumina results",
                            "Added task verbs to stop words; fetch 20 results (doc ranked #10 in shared sandbox)",
                            BOX_PURPLE),
        (13, "LATENCY",     "10–20s response time — 95% from Glean Chat API (sandbox)",
                            "fast_mode=True skips Chat → ~800ms; production fix is streaming",
                            BOX_ORANGE),
    ]

    for i, (num, cat, issue, fix, color) in enumerate(fixes):
        y = Inches(1.1 + i * 0.76)
        fix_row(slide, y, num, cat, issue, fix, color)

    # column headers
    for label, x in [("  #  CAT", Inches(0.25)), ("ISSUE", Inches(2.75)), ("FIX", Inches(7.7))]:
        add_text_box(slide, label, x, Inches(1.05), Inches(4.5), Inches(0.22),
                     font_size=8, bold=True, color=MID_GREY)

    # summary banner
    add_rect(slide, Inches(0.25), Inches(5.72), Inches(12.8), Inches(0.65), GLEAN_DARK)
    add_text_box(slide,
                 "13 fixes across indexing, search, Chat API, MCP transport, sandbox limitations, and UX. "
                 "Full detail in FIXES.md — github.com/karlhart/glean-chatbot",
                 Inches(0.45), Inches(5.8), Inches(12.4), Inches(0.5),
                 font_size=11, italic=True, color=WHITE, align=PP_ALIGN.CENTER)


def slide_4_files(prs):
    """Summary of key files."""
    slide = blank_slide(prs)
    fill_bg(slide, RGBColor(0xF8, 0xF9, 0xFB))

    # top bar
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.1), GLEAN_DARK)
    add_text_box(slide, "04", Inches(0.35), Inches(0.3), Inches(0.4), Inches(0.5),
                 font_size=10, bold=True, color=ACCENT_GOLD, align=PP_ALIGN.CENTER)
    add_text_box(slide, "KEY FILES",
                 Inches(0.85), Inches(0.28), Inches(5), Inches(0.45),
                 font_size=22, bold=True, color=WHITE)
    add_text_box(slide, "Repository: github.com/karlhart/glean-chatbot",
                 Inches(0.85), Inches(0.68), Inches(8), Inches(0.32),
                 font_size=12, color=MID_GREY)

    files = [
        {
            "name":  "src/indexer.py",
            "badge": "INDEXING API",
            "color": GLEAN_BLUE,
            "what":  "Pushes the 7 Lumina markdown documents into the Glean interviewds datasource using the Indexing API. Builds document payloads with stable IDs, plain-text body, viewURL, and open sandbox permissions. Supports both incremental (/indexdocuments) and bulk (/bulkindexdocuments) modes.",
            "why":   "Incremental mode chosen — re-running only updates changed docs without deleting the rest.",
        },
        {
            "name":  "src/chatbot.py",
            "badge": "SEARCH + CHAT APIs",
            "color": BOX_PURPLE,
            "what":  "Core RAG pipeline. Extracts keywords → fetches 20 results from Glean Search → filters by exact URL allowlist → enriches with full document content (snippet-first, 2500 char cap) → sends grounded prompt to Glean Chat → returns answer + sources. Includes fast_mode to skip Chat for ~800ms responses.",
            "why":   "Explicit context injection chosen over native Chat retrieval for reliable datasource scoping.",
        },
        {
            "name":  "src/mcp_server.py",
            "badge": "MCP TOOL",
            "color": BOX_GREEN,
            "what":  "FastMCP server exposing ask_lumina as a single callable tool for Claude Desktop. Parameters: question (required), datasource, top_k, include_citations, after_date, before_date, fast_mode. Instructions field tells Claude to present the full response including Sources verbatim.",
            "why":   "FastMCP auto-generates JSON schema from Python type hints — no manual OpenAPI spec needed.",
        },
        {
            "name":  "validate.py",
            "badge": "TEST SUITE",
            "color": BOX_ORANGE,
            "what":  "Five end-to-end test cases against the live Glean APIs. Checks: non-empty answer, at least one source returned, source titles match expected Lumina document keywords. All API calls log requestId and backendTimeMillis for debugging. Exit code 0 = all passed.",
            "why":   "Live API tests catch sandbox-specific issues (propagation lag, shared datasource pollution) that mocks would miss.",
        },
        {
            "name":  "docs/*.md  (7 files)",
            "badge": "DOCUMENT CORPUS",
            "color": RGBColor(0x55, 0x66, 0x88),
            "what":  "Employee Onboarding • IT Security • Production Workflow • Legal & Contracts • Post-Production / VFX • International Co-Production • Content Delivery Standards. Covers HR, IT, legal, production, and distribution — the full breadth of a media company's internal knowledge.",
            "why":   "Chosen to demonstrate Glean serving as a single source of truth across all departments.",
        },
    ]

    row_h = Inches(1.12)
    for i, f in enumerate(files):
        y = Inches(1.22 + i * row_h)

        # alternating row bg
        bg = RGBColor(0xFF, 0xFF, 0xFF) if i % 2 == 0 else LIGHT_GREY
        add_rect(slide, Inches(0.3), y, Inches(12.7), row_h - Inches(0.05), bg)

        # colour badge strip
        add_rect(slide, Inches(0.3), y, Inches(0.18), row_h - Inches(0.05), f["color"])

        # API badge pill
        add_rect(slide, Inches(0.6), y + Inches(0.12), Inches(1.8), Inches(0.28), f["color"])
        add_text_box(slide, f["badge"],
                     Inches(0.6), y + Inches(0.12), Inches(1.8), Inches(0.28),
                     font_size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # filename
        add_text_box(slide, f["name"],
                     Inches(2.55), y + Inches(0.08), Inches(3.5), Inches(0.35),
                     font_size=13, bold=True, color=GLEAN_DARK)

        # what it does
        add_text_box(slide, f["what"],
                     Inches(2.55), y + Inches(0.42), Inches(6.4), Inches(0.6),
                     font_size=9.5, color=RGBColor(0x44, 0x44, 0x55))

        # why label + text
        add_text_box(slide, "Why: " + f["why"],
                     Inches(9.1), y + Inches(0.12), Inches(3.8), Inches(0.85),
                     font_size=9, italic=True, color=RGBColor(0x66, 0x66, 0x77))

    # footer
    add_rect(slide, 0, Inches(6.95), SLIDE_W, Inches(0.55), GLEAN_DARK)
    add_text_box(slide,
                 "Built with Claude Code  •  Glean Indexing + Search + Chat APIs  •  Python / FastMCP",
                 Inches(0), Inches(7.0), SLIDE_W, Inches(0.4),
                 font_size=10, color=MID_GREY, align=PP_ALIGN.CENTER)


# ── Main ──────────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    prs = new_prs()
    slide_1_company(prs)
    slide_2_requirements(prs)
    slide_3_architecture(prs)
    slide_4_fixes_a(prs)
    slide_5_fixes_b(prs)

    out = "/Users/karl/Projects/GleanSADemo/LuminaChatbot_Presentation.pptx"
    prs.save(out)
    print(f"Saved → {out}")
